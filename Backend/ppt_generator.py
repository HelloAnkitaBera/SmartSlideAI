import os
import re
import json
import urllib.request
import urllib.parse
import tempfile
import google.generativeai as genai

from dotenv import load_dotenv


# =====================================
# IMAGE UTILS
# =====================================

def _download_topic_image(topic, width=480, height=1080):
    try:
        topic_lower = topic.lower()
        
        # Keywords categorization
        tech_keywords = ["data", "science", "code", "programming", "software", "ai", "machine", "cyber", "technology", "network", "cloud", "digital", "developer", "analytics", "computer", "system"]
        green_keywords = ["green", "sustain", "environment", "nature", "earth", "eco", "renewable", "solar", "wind", "climate", "clean", "energy"]
        biz_keywords = ["business", "finance", "corporate", "startup", "marketing", "sales", "strategy", "office", "money", "growth", "investment", "lead", "management", "project"]

        # Selection of curated high-quality vertical Unsplash photos
        if any(k in topic_lower for k in tech_keywords):
            urls = [
                "https://images.unsplash.com/photo-1635070041078-e363dbe005cb", # Sci-fi neon grid
                "https://images.unsplash.com/photo-1550751827-4bd374c3f58b", # Cybersecurity chip
                "https://images.unsplash.com/photo-1451187580459-43490279c0fa", # Cyber network
                "https://images.unsplash.com/photo-1526374965328-7f61d4dc18c5", # Coding matrix
            ]
        elif any(k in topic_lower for k in green_keywords):
            urls = [
                "https://images.unsplash.com/photo-1501854140801-50d01698950b", # Green forest top view
                "https://images.unsplash.com/photo-1447752875215-b2761acb3c5d", # Minimal forest path
            ]
        elif any(k in topic_lower for k in biz_keywords):
            urls = [
                "https://images.unsplash.com/photo-1486406146926-c627a92ad1ab", # Corporate architecture
                "https://images.unsplash.com/photo-1634017839464-5c339ebe3cb4", # 3D geometric shapes
                "https://images.unsplash.com/photo-1507679799987-c73779587ccf", # Premium suit minimal
            ]
        else:
            # Fallback abstract fluid art and tech lines
            urls = [
                "https://images.unsplash.com/photo-1618005182384-a83a8bd57fbe", # Abstract fluid wave
                "https://images.unsplash.com/photo-1618005198143-d528f99e3381", # Sleek tech network
                "https://images.unsplash.com/photo-1579546929518-9e396f3cc809", # Elegant gradient
            ]

        # Use a stable hash of the topic string to select one of the URLs consistently
        idx = sum(ord(c) for c in topic) % len(urls)
        base_url = urls[idx]
        
        # Append sizing & cropping parameters to the CDN URL
        url = f"{base_url}?w={width}&h={height}&fit=crop"
        
        req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
        with urllib.request.urlopen(req, timeout=8) as response:
            if response.status == 200:
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                temp_file.write(response.read())
                temp_file.close()
                return temp_file.name
    except Exception as e:
        print(f"Error downloading curated image: {e}")
    return None


from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


# =====================================
# GLOBAL CONFIG
# =====================================

FONT_NAME = "Arial"


# =====================================
# LOAD ENV
# =====================================

load_dotenv(override=True)

api_key = os.getenv("GEMINI_API_KEY")

if not api_key:
    raise ValueError("GEMINI_API_KEY missing in .env")

genai.configure(api_key=api_key)

# Using gemini-2.5-flash which is the active stable model
model = genai.GenerativeModel("gemini-2.5-flash")


# =====================================
# MODERN THEME PALETTES
# =====================================

THEMES = {
    "default": {
        "bg": RGBColor(8, 12, 28),
        "card": RGBColor(20, 25, 45),
        "primary": RGBColor(99, 102, 241),      # Indigo
        "secondary": RGBColor(168, 85, 247),    # Purple
        "text": RGBColor(255, 255, 255),
        "subtext": RGBColor(180, 180, 180),
    },
    "corporate": {
        "bg": RGBColor(15, 23, 42),             # Dark Slate Navy
        "card": RGBColor(30, 41, 59),            # Slate Card
        "primary": RGBColor(14, 165, 233),       # Sky Blue
        "secondary": RGBColor(234, 179, 8),      # Amber/Gold
        "text": RGBColor(255, 255, 255),
        "subtext": RGBColor(148, 163, 184),
    },
    "modern": {
        "bg": RGBColor(9, 9, 11),                # Zinc Black
        "card": RGBColor(24, 24, 27),            # Zinc Card
        "primary": RGBColor(16, 185, 129),       # Emerald
        "secondary": RGBColor(45, 212, 191),     # Teal
        "text": RGBColor(255, 255, 255),
        "subtext": RGBColor(161, 161, 170),
    },
    "startup": {
        "bg": RGBColor(17, 24, 39),              # Gray 900
        "card": RGBColor(31, 41, 55),            # Gray 800
        "primary": RGBColor(139, 92, 246),       # Violet
        "secondary": RGBColor(249, 115, 22),     # Orange
        "text": RGBColor(255, 255, 255),
        "subtext": RGBColor(156, 163, 175),
    }
}


# =====================================
# CLEAN TEXT (ASCII FIX)
# =====================================

def clean_text(text):
    if text is None:
        return ""
    text = str(text)
    # Keep safe ASCII characters
    text = text.encode("ascii", "ignore").decode("ascii")
    # Remove extra spaces
    text = re.sub(r"\s+", " ", text)
    return text.strip()


# =====================================
# FALLBACK CONTENT (DIVERSE LAYOUTS)
# =====================================

def fallback_content(topic, num_slides):
    slides = []
    layouts = ["standard", "two_column", "stat_highlight", "process", "summary"]
    
    for i in range(num_slides):
        layout = layouts[i % len(layouts)]
        
        slide_data = {
            "title": f"{topic}: Key Aspect {i+1}",
            "layout": layout,
            "paragraph": f"Analyzing the fundamental components of {topic} reveals key drivers for growth and innovation. This exploration delivers critical operational value.",
            "bullets": [
                f"Establishing core infrastructural drivers for long-term {topic} integration and development.",
                "Leveraging strategic design advantages for optimized deployment in modern workflows.",
                "Optimizing future scaling pathways with targeted metrics and continuous monitoring."
            ],
            "stat": "85%",
            "stat_label": "Accelerated Market Growth",
            "quote": f"{topic} represents the frontier of modern industry developments.",
            "image": f"Conceptual diagram showing {topic} integration and impact",
            "col1_title": "Current Implementation",
            "col1_bullets": [
                "Legacy infrastructure models present high maintenance overhead.",
                "Inflexible data architectures slow integration processes."
            ],
            "col2_title": "Future Integration",
            "col2_bullets": [
                "Scalable microservices patterns allow direct API interaction.",
                "Modernized cloud-native structures optimize processing speeds."
            ],
            "steps": [
                {"step_title": "Discovery & Planning", "step_desc": f"Assessing current bottlenecks of {topic} and establishing clear integration roadmaps."},
                {"step_title": "Implementation", "step_desc": "Launching key modules in staging environments with strict automated testing metrics."},
                {"step_title": "Optimization", "step_desc": "Monitoring performance targets and eliminating container scaling bottlenecks."}
            ]
        }
        slides.append(slide_data)
        
    return slides


# =====================================
# GENERATE CONTENT
# =====================================

def generate_content(topic, num_slides=5):
    prompt = f"""
Create EXACTLY {num_slides} professional presentation slides on:
"{topic}"

Each slide must focus on a DIFFERENT subtopic/concept related to "{topic}" and use a DIFFERENT layout for visual variety. Provide rich, detailed, and highly informative descriptions for each slide to make it a comprehensive guide.

Return a JSON array of slide objects.

Available slide layouts:
1. "standard": Best for a general concept. Requires 'paragraph', 'bullets', 'image', 'stat', 'stat_label', 'quote'.
2. "two_column": Best for comparison, options, or two distinct categories. Requires 'col1_title', 'col1_bullets', 'col2_title', 'col2_bullets', 'quote'.
3. "stat_highlight": Best for showing a massive statistic or key highlight. Requires 'stat', 'stat_label', 'paragraph', 'bullets', 'quote'.
4. "process": Best for workflow steps, timeline, phases, or progression. Requires 'steps' (array of 3 steps, each with 'step_title' and 'step_desc'), 'quote'.
5. "summary": Best for key takeaways. Requires 'paragraph', 'bullets', 'quote'.

JSON SCHEMA:
[
  {{
    "title": "Detailed slide title",
    "layout": "standard | two_column | stat_highlight | process | summary",
    "paragraph": "2-3 sentences of highly detailed, professional analysis providing rich context and explanation.",
    "bullets": [
      "Detail-rich bullet point offering a clear explanation or factual point (10-15 words)",
      "Detail-rich bullet point offering a clear explanation or factual point (10-15 words)",
      "Detail-rich bullet point offering a clear explanation or factual point (10-15 words)"
    ],
    "stat": "Key metric (e.g., '85%', '10x', '$4.2B')",
    "stat_label": "Descriptive metric label (3-5 words)",
    "quote": "A thoughtful quote, key takeaway, or strategic advice",
    "image": "Description of matching professional graphic, architectural diagram, or visual concept",
    "col1_title": "Title of Column 1",
    "col1_bullets": [
      "Detailed bullet explaining Column 1 aspect (10-15 words)",
      "Detailed bullet explaining Column 1 aspect (10-15 words)"
    ],
    "col2_title": "Title of Column 2",
    "col2_bullets": [
      "Detailed bullet explaining Column 2 aspect (10-15 words)",
      "Detailed bullet explaining Column 2 aspect (10-15 words)"
    ],
    "steps": [
      {{
        "step_title": "Detailed Title for Phase 1",
        "step_desc": "Thorough explanation of Phase 1 activities, requirements, and deliverables (8-12 words)"
      }},
      {{
        "step_title": "Detailed Title for Phase 2",
        "step_desc": "Thorough explanation of Phase 2 activities, requirements, and deliverables (8-12 words)"
      }},
      {{
        "step_title": "Detailed Title for Phase 3",
        "step_desc": "Thorough explanation of Phase 3 activities, requirements, and deliverables (8-12 words)"
      }}
    ]
  }}
]

Rules:
- Vary the layouts across slides (e.g. use standard, then process, then stat_highlight, then two_column, then summary).
- Do not repeat slide titles or contents.
- Keep all text professional, clear, and highly informative.
- Use only standard ASCII characters.
"""

    try:
        response = model.generate_content(
            prompt,
            generation_config={"response_mime_type": "application/json"}
        )

        text = response.text.strip()
        slides = json.loads(text)

        if not isinstance(slides, list):
            return fallback_content(topic, num_slides)

        cleaned_slides = []
        for slide in slides:
            cleaned_slide = {
                "title": clean_text(slide.get("title", "Slide")),
                "layout": clean_text(slide.get("layout", "standard")),
                "paragraph": clean_text(slide.get("paragraph", "")),
                "bullets": [clean_text(b) for b in slide.get("bullets", [])],
                "stat": clean_text(slide.get("stat", "")),
                "stat_label": clean_text(slide.get("stat_label", "")),
                "quote": clean_text(slide.get("quote", "")),
                "image": clean_text(slide.get("image", "")),
                "col1_title": clean_text(slide.get("col1_title", "Current")),
                "col1_bullets": [clean_text(b) for b in slide.get("col1_bullets", [])],
                "col2_title": clean_text(slide.get("col2_title", "Proposed")),
                "col2_bullets": [clean_text(b) for b in slide.get("col2_bullets", [])],
                "steps": []
            }
            
            steps = slide.get("steps", [])
            for step in steps:
                cleaned_slide["steps"].append({
                    "step_title": clean_text(step.get("step_title", "")),
                    "step_desc": clean_text(step.get("step_desc", ""))
                })
                
            cleaned_slides.append(cleaned_slide)

        return cleaned_slides

    except Exception as e:
        print("\nGemini Error:")
        print(e)
        return fallback_content(topic, num_slides)


# =====================================
# PPT GENERATOR
# =====================================

class PPTGenerator:

    def __init__(self, theme="default"):
        self.prs = Presentation()
        self.theme = THEMES.get(theme, THEMES["default"])

    def set_background(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["bg"]

    def add_card(self, slide, left, top, width, height, bg_color):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        # Border matches card background to be invisible/seamless
        card.line.color.rgb = bg_color
        return card

    def add_textbox(self, slide, left, top, width, height, word_wrap=True, margin=10):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = word_wrap
        tf.margin_left = Pt(margin)
        tf.margin_right = Pt(margin)
        tf.margin_top = Pt(margin)
        tf.margin_bottom = Pt(margin)
        return tf

    def add_slide_title(self, slide, title):
        title_box = slide.shapes.add_textbox(Pt(40), Pt(20), Pt(640), Pt(50))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(0)
        p = tf.paragraphs[0]
        p.text = clean_text(title)
        p.font.name = FONT_NAME
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.theme["text"]

    def add_slide_quote(self, slide, quote):
        if not quote:
            return
        quote_box = slide.shapes.add_textbox(Pt(40), Pt(425), Pt(640), Pt(40))
        tf = quote_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(0)
        p = tf.paragraphs[0]
        p.text = f'"{clean_text(quote)}"'
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = self.theme["subtext"]

    # =================================
    # TITLE SLIDE (COVER)
    # =================================

    def add_title_slide(self, topic):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_background(slide)

        # Try to download and add topic-relevant cover image
        image_path = _download_topic_image(topic, width=480, height=1080)
        image_added = False
        if image_path:
            try:
                slide.shapes.add_picture(image_path, Pt(480), Pt(0), Pt(240), Pt(540))
                image_added = True
            except Exception as e:
                print(f"Failed to add title slide picture: {e}")
            finally:
                try:
                    os.remove(image_path)
                except:
                    pass

        # Fallback to elegant vertical block accent and badge if image download failed
        if not image_added:
            self.add_card(
                slide,
                Pt(480), Pt(0), Pt(240), Pt(540),
                bg_color=self.theme["primary"]
            )

            # Decorative secondary accent badge
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Pt(550), Pt(210), Pt(100), Pt(100)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme["secondary"]
            circle.line.color.rgb = self.theme["secondary"]

        # Small badge label
        tf_badge = self.add_textbox(slide, Pt(60), Pt(120), Pt(400), Pt(40))
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = "PRESENTATION TOPIC"
        p_badge.font.name = FONT_NAME
        p_badge.font.size = Pt(12)
        p_badge.font.bold = True
        p_badge.font.color.rgb = self.theme["primary"]

        # Title
        tf_title = self.add_textbox(slide, Pt(55), Pt(150), Pt(400), Pt(180))
        p_title = tf_title.paragraphs[0]
        p_title.text = clean_text(topic)
        p_title.font.name = FONT_NAME
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = self.theme["text"]

        # Accent Line
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Pt(60), Pt(340), Pt(150), Pt(4)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["secondary"]
        line.line.color.rgb = self.theme["secondary"]

        # Subtitle
        tf_sub = self.add_textbox(slide, Pt(55), Pt(360), Pt(400), Pt(80))
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = "AI Generated Professional Presentation"
        p_sub.font.name = FONT_NAME
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = self.theme["subtext"]

    # =================================
    # LAYOUT 1: STANDARD SLIDE
    # =================================

    def add_standard_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_background(slide)
        self.add_slide_title(slide, data.get("title", "Overview"))

        # Left Column Card (Main details)
        self.add_card(slide, Pt(40), Pt(90), Pt(430), Pt(320), bg_color=self.theme["card"])
        
        # Paragraph (2-3 sentences, adjusted padding & size to fit)
        tf_para = self.add_textbox(slide, Pt(50), Pt(100), Pt(410), Pt(80))
        p_para = tf_para.paragraphs[0]
        p_para.text = clean_text(data.get("paragraph", ""))
        p_para.font.name = FONT_NAME
        p_para.font.size = Pt(11.5)
        p_para.font.color.rgb = self.theme["subtext"]

        # Bullets (Detailed bullets, adjusted font size to fit)
        tf_bullets = self.add_textbox(slide, Pt(50), Pt(185), Pt(410), Pt(210))
        bullets = data.get("bullets", [])
        for i, bullet in enumerate(bullets):
            p = tf_bullets.paragraphs[0] if i == 0 else tf_bullets.add_paragraph()
            p.text = f"•  {clean_text(bullet)}"
            p.font.name = FONT_NAME
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme["text"]
            p.space_after = Pt(6)

        # Right Column Card (Visual context block)
        self.add_card(slide, Pt(495), Pt(90), Pt(185), Pt(190), bg_color=self.theme["primary"])
        
        tf_visual = self.add_textbox(slide, Pt(505), Pt(110), Pt(165), Pt(150))
        p_vis = tf_visual.paragraphs[0]
        p_vis.text = clean_text(data.get("image", "Visual Concept"))
        p_vis.font.name = FONT_NAME
        p_vis.font.size = Pt(13)
        p_vis.font.bold = True
        p_vis.font.color.rgb = RGBColor(255, 255, 255)
        p_vis.alignment = 1

        # Stat Callout
        if data.get("stat"):
            tf_stat = self.add_textbox(slide, Pt(495), Pt(295), Pt(185), Pt(110))
            p_stat = tf_stat.paragraphs[0]
            p_stat.text = clean_text(data.get("stat"))
            p_stat.font.name = FONT_NAME
            p_stat.font.size = Pt(28)
            p_stat.font.bold = True
            p_stat.font.color.rgb = self.theme["secondary"]
            p_stat.alignment = 1
            
            p_label = tf_stat.add_paragraph()
            p_label.text = clean_text(data.get("stat_label", "Key Metric"))
            p_label.font.name = FONT_NAME
            p_label.font.size = Pt(11)
            p_label.font.color.rgb = self.theme["subtext"]
            p_label.alignment = 1

        self.add_slide_quote(slide, data.get("quote", ""))

    # =================================
    # LAYOUT 2: TWO COLUMN SLIDE
    # =================================

    def add_two_column_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_background(slide)
        self.add_slide_title(slide, data.get("title", "Comparison Analysis"))

        # Left Column Card
        self.add_card(slide, Pt(40), Pt(90), Pt(305), Pt(320), bg_color=self.theme["card"])
        col1_title = data.get("col1_title", "Current State")
        tf_col1 = self.add_textbox(slide, Pt(50), Pt(100), Pt(285), Pt(300))
        p1 = tf_col1.paragraphs[0]
        p1.text = clean_text(col1_title)
        p1.font.name = FONT_NAME
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = self.theme["primary"]
        p1.space_after = Pt(10)

        col1_bullets = data.get("col1_bullets", [])
        if not col1_bullets:
            col1_bullets = data.get("bullets", [])[:2]
        
        for bullet in col1_bullets:
            p = tf_col1.add_paragraph()
            p.text = f"•  {clean_text(bullet)}"
            p.font.name = FONT_NAME
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme["text"]
            p.space_after = Pt(6)

        # Right Column Card
        self.add_card(slide, Pt(375), Pt(90), Pt(305), Pt(320), bg_color=self.theme["card"])
        col2_title = data.get("col2_title", "Proposed Strategy")
        tf_col2 = self.add_textbox(slide, Pt(385), Pt(100), Pt(285), Pt(300))
        p2 = tf_col2.paragraphs[0]
        p2.text = clean_text(col2_title)
        p2.font.name = FONT_NAME
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = self.theme["secondary"]
        p2.space_after = Pt(10)

        col2_bullets = data.get("col2_bullets", [])
        if not col2_bullets:
            col2_bullets = data.get("bullets", [])[2:] if len(data.get("bullets", [])) > 2 else ["Optimized processes", "Scalable outcomes"]
        
        for bullet in col2_bullets:
            p = tf_col2.add_paragraph()
            p.text = f"•  {clean_text(bullet)}"
            p.font.name = FONT_NAME
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme["text"]
            p.space_after = Pt(6)

        self.add_slide_quote(slide, data.get("quote", ""))

    # =================================
    # LAYOUT 3: STAT HIGHLIGHT SLIDE
    # =================================

    def add_stat_highlight_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_background(slide)
        self.add_slide_title(slide, data.get("title", "Impact Evaluation"))

        # Left Column Card (Highlight colored)
        self.add_card(slide, Pt(40), Pt(90), Pt(280), Pt(320), bg_color=self.theme["primary"])
        
        # Stat Number
        stat_val = data.get("stat", "95%")
        tf_stat = self.add_textbox(slide, Pt(50), Pt(120), Pt(260), Pt(260))
        p_stat = tf_stat.paragraphs[0]
        p_stat.text = clean_text(stat_val)
        p_stat.font.name = FONT_NAME
        p_stat.font.size = Pt(56)
        p_stat.font.bold = True
        p_stat.font.color.rgb = RGBColor(255, 255, 255)
        p_stat.alignment = 1
        p_stat.space_after = Pt(12)

        # Stat Label
        p_label = tf_stat.add_paragraph()
        p_label.text = clean_text(data.get("stat_label", "Growth Index"))
        p_label.font.name = FONT_NAME
        p_label.font.size = Pt(14)
        p_label.font.bold = True
        p_label.font.color.rgb = self.theme["secondary"]
        p_label.alignment = 1

        # Right Column Card (Supporting content)
        self.add_card(slide, Pt(340), Pt(90), Pt(340), Pt(320), bg_color=self.theme["card"])
        
        tf_content = self.add_textbox(slide, Pt(350), Pt(100), Pt(320), Pt(300))
        p_title = tf_content.paragraphs[0]
        p_title.text = "Key Metrics & Support"
        p_title.font.name = FONT_NAME
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = self.theme["text"]
        p_title.space_after = Pt(10)

        p_para = tf_content.add_paragraph()
        p_para.text = clean_text(data.get("paragraph", ""))
        p_para.font.name = FONT_NAME
        p_para.font.size = Pt(11.5)
        p_para.font.color.rgb = self.theme["subtext"]
        p_para.space_after = Pt(10)

        bullets = data.get("bullets", [])
        for bullet in bullets:
            p_b = tf_content.add_paragraph()
            p_b.text = f"•  {clean_text(bullet)}"
            p_b.font.name = FONT_NAME
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = self.theme["text"]
            p_b.space_after = Pt(5)

        self.add_slide_quote(slide, data.get("quote", ""))

    # =================================
    # LAYOUT 4: PROCESS / TIMELINE SLIDE
    # =================================

    def add_process_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_background(slide)
        self.add_slide_title(slide, data.get("title", "Phased Execution Plan"))

        steps = data.get("steps", [])
        if not steps or len(steps) < 3:
            steps = [
                {"step_title": "Planning", "step_desc": "Define criteria and map current workflows."},
                {"step_title": "Deployment", "step_desc": "Release components inside sandbox clusters."},
                {"step_title": "Review", "step_desc": "Measure outcome key metrics against baselines."}
            ]

        # Draw 3 horizontal steps cards
        card_width = Pt(200)
        card_height = Pt(300)
        y_pos = Pt(100)
        
        for i, step in enumerate(steps[:3]):
            x_pos = Pt(40 + i * 220)
            
            # Step Card
            self.add_card(slide, x_pos, y_pos, card_width, card_height, bg_color=self.theme["card"])
            
            # Step Content
            tf_step = self.add_textbox(slide, x_pos + Pt(10), y_pos + Pt(10), card_width - Pt(20), card_height - Pt(20))
            
            # Step Index
            p_num = tf_step.paragraphs[0]
            p_num.text = f"0{i+1}"
            p_num.font.name = FONT_NAME
            p_num.font.size = Pt(18)
            p_num.font.bold = True
            p_num.font.color.rgb = self.theme["secondary"]
            p_num.space_after = Pt(8)
            
            # Step Title
            p_title = tf_step.add_paragraph()
            p_title.text = clean_text(step.get("step_title", f"Phase {i+1}"))
            p_title.font.name = FONT_NAME
            p_title.font.size = Pt(13)
            p_title.font.bold = True
            p_title.font.color.rgb = self.theme["primary"]
            p_title.space_after = Pt(8)
            
            # Step Desc
            p_desc = tf_step.add_paragraph()
            p_desc.text = clean_text(step.get("step_desc", ""))
            p_desc.font.name = FONT_NAME
            p_desc.font.size = Pt(10.5)
            p_desc.font.color.rgb = self.theme["subtext"]

        self.add_slide_quote(slide, data.get("quote", ""))

    # =================================
    # LAYOUT 5: SUMMARY SLIDE
    # =================================

    def add_summary_slide(self, data):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_background(slide)
        self.add_slide_title(slide, data.get("title", "Key Takeaways"))

        # Left Card (Key findings)
        self.add_card(slide, Pt(40), Pt(90), Pt(400), Pt(320), bg_color=self.theme["card"])
        tf_left = self.add_textbox(slide, Pt(50), Pt(100), Pt(380), Pt(300))
        
        p_hdr = tf_left.paragraphs[0]
        p_hdr.text = "Executive Summary"
        p_hdr.font.name = FONT_NAME
        p_hdr.font.size = Pt(15)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = self.theme["primary"]
        p_hdr.space_after = Pt(10)

        p_para = tf_left.add_paragraph()
        p_para.text = clean_text(data.get("paragraph", "Key points of implementation summary."))
        p_para.font.name = FONT_NAME
        p_para.font.size = Pt(11.5)
        p_para.font.color.rgb = self.theme["subtext"]
        p_para.space_after = Pt(10)

        bullets = data.get("bullets", [])
        for bullet in bullets:
            p_b = tf_left.add_paragraph()
            p_b.text = f"✔  {clean_text(bullet)}"
            p_b.font.name = FONT_NAME
            p_b.font.size = Pt(12)
            p_b.font.color.rgb = self.theme["text"]
            p_b.space_after = Pt(6)

        # Right Card (Call to Action)
        self.add_card(slide, Pt(465), Pt(90), Pt(215), Pt(320), bg_color=self.theme["secondary"])
        tf_right = self.add_textbox(slide, Pt(475), Pt(120), Pt(195), Pt(260))
        
        p_cta = tf_right.paragraphs[0]
        p_cta.text = "Next Steps"
        p_cta.font.name = FONT_NAME
        p_cta.font.size = Pt(18)
        p_cta.font.bold = True
        p_cta.font.color.rgb = RGBColor(255, 255, 255)
        p_cta.alignment = 1
        p_cta.space_after = Pt(15)

        p_cta_desc = tf_right.add_paragraph()
        p_cta_desc.text = clean_text(data.get("quote", "Execute findings to accelerate delivery."))
        p_cta_desc.font.name = FONT_NAME
        p_cta_desc.font.size = Pt(12)
        p_cta_desc.font.italic = True
        p_cta_desc.font.color.rgb = RGBColor(255, 255, 255)
        p_cta_desc.alignment = 1

    # =================================
    # LAYOUT 6: THANK YOU SLIDE
    # =================================

    def add_thank_you_slide(self, topic):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.set_background(slide)

        # Try to download and add topic-relevant conclusion image
        image_path = _download_topic_image(topic, width=480, height=1080)
        image_added = False
        if image_path:
            try:
                slide.shapes.add_picture(image_path, Pt(0), Pt(0), Pt(240), Pt(540))
                image_added = True
            except Exception as e:
                print(f"Failed to add thank you slide picture: {e}")
            finally:
                try:
                    os.remove(image_path)
                except:
                    pass

        # Fallback to symmetrical visual block and badge if image download failed
        if not image_added:
            self.add_card(
                slide,
                Pt(0), Pt(0), Pt(240), Pt(540),
                bg_color=self.theme["primary"]
            )

            # Decorative secondary accent badge
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Pt(70), Pt(210), Pt(100), Pt(100)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme["secondary"]
            circle.line.color.rgb = self.theme["secondary"]

        # Small badge label
        tf_badge = self.add_textbox(slide, Pt(280), Pt(160), Pt(400), Pt(40))
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = "CONCLUSION"
        p_badge.font.name = FONT_NAME
        p_badge.font.size = Pt(12)
        p_badge.font.bold = True
        p_badge.font.color.rgb = self.theme["primary"]

        # Main Thank You title
        tf_title = self.add_textbox(slide, Pt(275), Pt(190), Pt(400), Pt(120))
        p_title = tf_title.paragraphs[0]
        p_title.text = "Thank You"
        p_title.font.name = FONT_NAME
        p_title.font.size = Pt(48)
        p_title.font.bold = True
        p_title.font.color.rgb = self.theme["text"]

        # Accent Line
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Pt(280), Pt(320), Pt(150), Pt(4)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["secondary"]
        line.line.color.rgb = self.theme["secondary"]

        # Subtitle / Closing statement
        tf_sub = self.add_textbox(slide, Pt(275), Pt(340), Pt(400), Pt(100))
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = f"We hope this overview of {clean_text(topic)} was insightful. Let's connect for any questions or next steps."
        p_sub.font.name = FONT_NAME
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = self.theme["subtext"]

    # =================================
    # SAVE PPT
    # =================================

    def save(self, topic):
        output_dir = os.path.abspath(
            os.path.join(
                os.path.dirname(__file__),
                "..",
                "Frontend",
                "generated"
            )
        )
        os.makedirs(output_dir, exist_ok=True)

        safe_name = "".join(c for c in topic if c.isalnum() or c in (" ", "_"))
        filename = safe_name.replace(" ", "_") + ".pptx"
        path = os.path.join(output_dir, filename)

        if os.path.exists(path):
            os.remove(path)

        self.prs.save(path)

        # Validate PPT works
        Presentation(path)
        size = os.path.getsize(path)

        print("\nPPT SUCCESSFULLY SAVED")
        print(f"SIZE: {size} bytes")
        print(f"LOCATION: {path}")

        return path

    # =================================
    # CREATE PRESENTATION
    # =================================

    def create_presentation(self, topic, num_slides=5):
        # 1. Add Title Cover Slide
        self.add_title_slide(topic)

        # 2. Call Gemini API to generate structured content for remaining slides
        # Slide 1 is Title, Last slide is Thank You.
        content_slides_count = max(1, num_slides - 2)
        slides_data = generate_content(topic, content_slides_count)

        # 3. Route slides data to the specified layouts
        for slide_data in slides_data:
            layout_type = slide_data.get("layout", "standard").strip().lower()
            
            if layout_type == "two_column":
                self.add_two_column_slide(slide_data)
            elif layout_type == "stat_highlight":
                self.add_stat_highlight_slide(slide_data)
            elif layout_type == "process":
                self.add_process_slide(slide_data)
            elif layout_type == "summary":
                self.add_summary_slide(slide_data)
            else:
                # default fallback layout
                self.add_standard_slide(slide_data)

        # 4. Add Thank You Slide
        self.add_thank_you_slide(topic)

        # 5. Save presentation to file
        return self.save(topic)


# =====================================
# MAIN EXPOSED FUNCTION
# =====================================

def create_ppt(topic, num_slides=5, theme="default"):
    ppt = PPTGenerator(theme)
    return ppt.create_presentation(topic, num_slides)


# =====================================
# TEST RUN
# =====================================

if __name__ == "__main__":
    try:
        ppt_path = create_ppt(
            topic="Artificial Intelligence",
            num_slides=6,
            theme="default"
        )
        print("\nPPT CREATED SUCCESSFULLY")
        print(ppt_path)
    except Exception as e:
        print("\nERROR:")
        print(e)